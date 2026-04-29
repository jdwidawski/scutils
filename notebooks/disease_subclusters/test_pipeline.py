"""
End-to-end smoke test for the disease subcluster pipeline.

Creates a synthetic AnnData that mimics real data (disease/cell-type labels,
kNN graph, UMAP), then runs Notebook 01 (detection + DE) and Notebook 02
(HTML/PPTX card generation) programmatically.

Run with:
    uv run python notebooks/disease_subclusters/test_pipeline.py
"""
from __future__ import annotations

import re
import shutil
import sys
import tempfile
from datetime import datetime
from pathlib import Path

import matplotlib
matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import scanpy as sc
from anndata import AnnData
from scipy.sparse import csr_matrix

# ── Make sure the package is importable ───────────────────────────────────────
sys.path.insert(0, str(Path(__file__).parents[2] / "src"))

import scutils
from scutils.tools.disease_subclusters import (
    detect_disease_enriched_subclusters,
    plot_disease_enriched_subclusters,
)
from scutils.tools.differential_expression import deseq2, format_deseq2_results
from scutils.preprocessing.pseudobulk import pseudobulk
from scutils.plotting.volcano_plot import volcano_plot

# ─────────────────────────────────────────────────────────────────────────────
# 1. Synthetic dataset
# ─────────────────────────────────────────────────────────────────────────────

def make_synthetic_adata(
    n_healthy: int = 800,
    n_disease: int = 600,
    n_genes: int = 200,
    n_samples: int = 8,
    seed: int = 42,
) -> AnnData:
    """Return a synthetic AnnData with two diseases, multiple cell types,
    raw count layer, UMAP embedding, and kNN graph."""
    rng = np.random.default_rng(seed)
    n_cells = n_healthy + n_disease

    # ── Expression matrix ─────────────────────────────────────────────────────
    X = rng.poisson(2.0, size=(n_cells, n_genes)).astype(np.float32)

    # Inject a disease signature: first 30 genes upregulated in disease cells
    X[n_healthy:, :30] += rng.poisson(5.0, size=(n_disease, 30)).astype(np.float32)

    # ── obs metadata ─────────────────────────────────────────────────────────
    cell_types = np.array(
        ["Epithelial"] * (n_healthy // 2)
        + ["Macrophage"] * (n_healthy // 2)
        + ["Epithelial"] * (n_disease // 2)
        + ["Macrophage"] * (n_disease // 2)
    )
    diseases = np.array(
        ["Control"] * n_healthy
        + ["DiseaseA"] * (n_disease // 2)
        + ["DiseaseB"] * (n_disease // 2)
    )
    # Assign samples (biological replicates), balanced across conditions
    samples = np.array(
        [f"S{(i % (n_samples // 2)) + 1}" for i in range(n_healthy)]   # S1–S4 healthy
        + [f"S{(i % (n_samples // 2)) + n_samples // 2 + 1}" for i in range(n_disease)]  # S5–S8 disease
    )

    obs = pd.DataFrame(
        {
            "disease":   diseases,
            "cell_type": cell_types,
            "sampleID":  samples,
        },
        index=[f"cell_{i}" for i in range(n_cells)],
    )

    var = pd.DataFrame(
        {"gene_symbol": [f"Gene{i}" for i in range(n_genes)]},
        index=[f"Gene{i}" for i in range(n_genes)],
    )

    adata = AnnData(X=csr_matrix(X), obs=obs, var=var)
    adata.layers["counts"] = adata.X.copy()

    # ── PCA + kNN + UMAP ─────────────────────────────────────────────────────
    sc.pp.normalize_total(adata, target_sum=1e4)
    sc.pp.log1p(adata)
    sc.pp.pca(adata, n_comps=20)
    sc.pp.neighbors(adata, n_neighbors=15, use_rep="X_pca")
    sc.tl.umap(adata)

    return adata


# ─────────────────────────────────────────────────────────────────────────────
# 2. NB01 — Detection + DE
# ─────────────────────────────────────────────────────────────────────────────

def run_nb01(adata: AnnData, out_dir: Path) -> pd.DataFrame:
    print("\n" + "=" * 60)
    print("NB01: Disease subcluster detection + DE")
    print("=" * 60)

    # Detection
    detect_disease_enriched_subclusters(
        adata,
        disease_key="disease",
        celltype_key="cell_type",
        groups_disease=["DiseaseA", "DiseaseB"],
        reference_group=["Control"],
        combine_diseases="pool",
        min_enrichment_fold=1.2,
        min_subcluster_size=30,
        enrichment_fdr=0.1,
        spatial_sensitivity="low",
        min_reference_cells=20,
        result_key="disease_subcluster",
        verbose=False,
    )

    labels = adata.obs["disease_subcluster"]
    subclusters = sorted(labels[labels != "background"].unique())
    print(f"  Detected {len(subclusters)} subcluster(s): {subclusters}")

    info_df = adata.uns.get("disease_subcluster_info", pd.DataFrame())
    print(f"  Info table rows: {len(info_df)}")

    # Save info CSV
    info_path = out_dir / "test_dataset_subcluster_info.csv"
    info_df.to_csv(info_path, index=False)
    print(f"  Saved: {info_path}")

    # Save AnnData
    adata_path = out_dir / "test_dataset_disease_subclusters.h5ad"
    adata.write_h5ad(adata_path)
    print(f"  Saved: {adata_path}")

    # Visualise (smoke test)
    fig = plot_disease_enriched_subclusters(
        adata, celltype_key="cell_type", disease_key="disease",
        result_key="disease_subcluster", split_by="disease", show=False,
    )
    fig.savefig(out_dir / "subclusters_by_disease.png", bbox_inches="tight", dpi=80)
    plt.close(fig)
    print("  Visualisation saved.")

    # DE with pyDESeq2
    de_dir = out_dir / "de"
    de_dir.mkdir(exist_ok=True)

    de_results: dict[str, pd.DataFrame] = {}

    for sc_label in subclusters:
        test_mask = labels == sc_label
        ref_mask  = ~test_mask
        test_n = test_mask.sum()
        ref_n  = ref_mask.sum()

        # Check replicates
        def _n_replicates(mask, min_cells=5):
            return (
                adata.obs.loc[mask]
                .groupby("sampleID")
                .size()
                .pipe(lambda s: (s >= min_cells).sum())
            )

        n_test_reps = _n_replicates(test_mask)
        n_ref_reps  = _n_replicates(ref_mask)

        if n_test_reps < 2 or n_ref_reps < 2:
            print(f"  ⚠  {sc_label}: insufficient replicates "
                  f"(test={n_test_reps}, ref={n_ref_reps}) — skipping DE")
            continue

        print(f"  Running DE: {sc_label} ({test_n} test, {ref_n} ref cells)", end=" ... ")
        try:
            adata_sub = adata[test_mask | ref_mask].copy()
            adata_sub.obs["_group"] = "ref"
            adata_sub.obs.loc[adata_sub.obs_names.isin(adata.obs_names[test_mask]), "_group"] = "test"

            pb = pseudobulk(
                adata_sub,
                sample_col="sampleID",
                groups_col="_group",
                layer="counts",
                min_cells=5,
            )

            res = deseq2(
                pb,
                design="~_group",
                contrast=["_group", "test", "ref"],
                alpha=0.05,
                shrink_lfc=True,
                quiet=True,
            )
            res["gene"] = res.index
            res["comparison"] = f"{sc_label}_vs_rest"
            res.reset_index(drop=True, inplace=True)

            safe = re.sub(r'[^\w\-]', '_', sc_label).strip('_')
            csv_path = de_dir / f"de_{safe}.csv"
            res.to_csv(csv_path, index=False)
            de_results[sc_label] = res
            print(f"OK ({len(res)} genes)")
        except Exception as e:
            print(f"FAILED: {e}")

    print(f"\n  DE complete: {len(de_results)}/{len(subclusters)} subclusters")
    return info_df


# ─────────────────────────────────────────────────────────────────────────────
# 3. NB02 — Subcluster cards (HTML + PPTX)
# ─────────────────────────────────────────────────────────────────────────────

def run_nb02(adata: AnnData, results_dir: Path, tissue: str = "test_tissue") -> None:
    print("\n" + "=" * 60)
    print("NB02: Subcluster cards (HTML + PPTX)")
    print("=" * 60)

    # Local imports matching NB02
    import io, re as _re, shutil
    from jinja2 import Environment, FileSystemLoader
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Cm, Inches, Pt

    RESULTS_DIR     = results_dir
    TISSUE          = tissue
    SUBCLUSTER_KEY  = "disease_subcluster"
    CELLTYPE_KEY    = "cell_type"
    DISEASE_KEY     = "disease"
    DE_PVAL_COL     = "padj"
    DE_LFC_COL      = "log2FoldChange"
    GENE_SYMBOL_COL = "gene"
    UMAP_FIGSIZE    = (5, 4)
    UMAP_POINT_SIZE = None
    BACKGROUND_COLOR = "#d3d3d3"
    VOLCANO_PVAL_CUTOFF = 0.05
    VOLCANO_LFC_CUTOFF  = 0.5
    VOLCANO_TOP_N_UP    = 5
    VOLCANO_TOP_N_DOWN  = 3
    VOLCANO_FIGSIZE     = (8, 5)

    REPORT_DIR     = RESULTS_DIR / f"{TISSUE}_subcluster_report"
    REPORT_IMG_DIR = REPORT_DIR / "img"
    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    REPORT_IMG_DIR.mkdir(exist_ok=True)
    (REPORT_DIR / "de").mkdir(exist_ok=True)

    def _fig_to_bytes(fig):
        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", dpi=80)
        buf.seek(0)
        return buf.read()

    def _safe_label(label):
        safe = _re.sub(r'[^\w\-]', '_', label)
        safe = _re.sub(r'_+', '_', safe)
        return safe.strip('_')

    def _make_umap_celltype(adata, figsize):
        fig, ax = plt.subplots(figsize=figsize)
        sc.pl.umap(adata, color=CELLTYPE_KEY, ax=ax, show=False, frameon=False)
        return fig

    def _make_umap_disease(adata, figsize):
        fig, ax = plt.subplots(figsize=figsize)
        sc.pl.umap(adata, color=DISEASE_KEY, ax=ax, show=False, frameon=False)
        return fig

    def _make_umap_highlight(adata, subcluster_label, figsize):
        mask = adata.obs[SUBCLUSTER_KEY] == subcluster_label
        fig, ax = plt.subplots(figsize=figsize)
        coords = adata.obsm["X_umap"]
        pt_size = UMAP_POINT_SIZE or max(120000 / adata.n_obs, 0.5)
        ax.scatter(coords[~mask, 0], coords[~mask, 1], s=pt_size, c=BACKGROUND_COLOR, linewidths=0)
        ax.scatter(coords[mask, 0], coords[mask, 1], s=pt_size * 1.5, c="#e63946", linewidths=0, zorder=2)
        ax.axis("off")
        return fig

    # Load subcluster info
    info_path = RESULTS_DIR / f"test_dataset_subcluster_info.csv"
    subcluster_info = pd.read_csv(info_path) if info_path.exists() else pd.DataFrame()

    labels = adata.obs[SUBCLUSTER_KEY]
    subcluster_labels = sorted(labels[labels != "background"].unique())
    print(f"  Building cards for {len(subcluster_labels)} subclusters")

    subcluster_data = []

    for label in subcluster_labels:
        safe = _safe_label(label)
        subcluster_img_dir = REPORT_IMG_DIR / safe
        subcluster_img_dir.mkdir(exist_ok=True)

        card = {
            "label":  label,
            "anchor": safe,
            "disease": label.split("|")[1] if "|" in label else label,
            "stats": {},
            "disease_breakdown": [],
        }

        # UMAPs
        for key, make_fn, fname in [
            ("umap_celltype",  _make_umap_celltype,  "umap_celltype.png"),
            ("umap_disease",   _make_umap_disease,   "umap_disease.png"),
            ("umap_highlight", lambda a, fs: _make_umap_highlight(a, label, fs), "umap_highlight.png"),
        ]:
            try:
                fig = make_fn(adata, UMAP_FIGSIZE)
                fig.savefig(subcluster_img_dir / fname, bbox_inches="tight", dpi=80)
                card[f"{key}_path"]  = f"img/{safe}/{fname}"
                card[f"{key}_bytes"] = _fig_to_bytes(fig)
                plt.close(fig)
            except Exception as e:
                print(f"  [warn] {key} failed: {e}")
                card[f"{key}_path"] = card[f"{key}_bytes"] = None

        # DE
        de_df = None
        de_dir = RESULTS_DIR / "de"
        if de_dir.exists():
            candidates = [p for p in sorted(de_dir.glob(f"de_{safe}*.csv")) if "all_comparisons" not in p.name]
            if candidates:
                de_df = pd.read_csv(candidates[0])

        # Copy DE CSV to report
        if de_df is not None:
            dst = REPORT_DIR / "de" / f"{safe}.csv"
            de_df.to_csv(dst, index=False)
            card["de_csv_path"] = f"de/{safe}.csv"
            card["de_json"]     = de_df.to_json(orient="split")
        else:
            card["de_csv_path"] = None
            card["de_json"]     = None

        # Volcano
        if de_df is not None and GENE_SYMBOL_COL in de_df.columns:
            try:
                df_v = format_deseq2_results(de_df.set_index(GENE_SYMBOL_COL), pval_col=DE_PVAL_COL, lfc_col=DE_LFC_COL)
                fig_v = volcano_plot(df_v, pval_cutoff=VOLCANO_PVAL_CUTOFF, lfc_cutoff=VOLCANO_LFC_CUTOFF,
                                     top_n_up=VOLCANO_TOP_N_UP, top_n_down=VOLCANO_TOP_N_DOWN, figsize=VOLCANO_FIGSIZE)
                fig_v.savefig(subcluster_img_dir / "volcano.png", bbox_inches="tight", dpi=80)
                card["volcano_path"]  = f"img/{safe}/volcano.png"
                card["volcano_bytes"] = _fig_to_bytes(fig_v)
                plt.close(fig_v)
            except Exception as e:
                print(f"  [warn] volcano failed: {e}")
                card["volcano_path"] = card["volcano_bytes"] = None
        else:
            card["volcano_path"] = card["volcano_bytes"] = None

        card["dotplot_path"]  = None
        card["dotplot_bytes"] = None

        subcluster_data.append(card)
        print(f"  ✓ {label}")

    # ── Render HTML ───────────────────────────────────────────────────────────
    template_dir = Path(__file__).parent / "templates"
    env = Environment(loader=FileSystemLoader(str(template_dir)), autoescape=False)
    template = env.get_template("subcluster_report.html.j2")
    html = template.render(
        tissue=TISSUE,
        subclusters=subcluster_data,
        generated_at=datetime.now().strftime("%Y-%m-%d %H:%M"),
    )
    html_path = REPORT_DIR / "index.html"
    html_path.write_text(html, encoding="utf-8")
    print(f"\n  ✓ HTML written: {html_path} ({html_path.stat().st_size // 1024} KB)")

    # ── Build PPTX ───────────────────────────────────────────────────────────
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for card in subcluster_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Cm(1), Cm(0.3), Cm(28), Cm(1.5))
        txBox.text_frame.text = card["label"]

    pptx_path = REPORT_DIR / f"{TISSUE}_subcluster_cards.pptx"
    prs.save(str(pptx_path))
    print(f"  ✓ PPTX written: {pptx_path} ({pptx_path.stat().st_size // 1024} KB)")


# ─────────────────────────────────────────────────────────────────────────────
# 4. NB03 — Combined report (HTML)
# ─────────────────────────────────────────────────────────────────────────────

def run_nb03_combined(results_dir: Path) -> None:
    print("\n" + "=" * 60)
    print("NB03: Combined report")
    print("=" * 60)

    import re, shutil
    from jinja2 import Environment, FileSystemLoader
    from pptx import Presentation
    from pptx.util import Inches, Cm

    TISSUE_REPORT_DIRS = {"test_tissue": results_dir / "test_tissue_subcluster_report"}
    OUTPUT_DIR = results_dir / "combined_report"
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    (OUTPUT_DIR / "img").mkdir(exist_ok=True)
    (OUTPUT_DIR / "de").mkdir(exist_ok=True)

    tissues: dict[str, list[dict]] = {}

    for tissue, report_dir in TISSUE_REPORT_DIRS.items():
        if not report_dir.exists():
            print(f"  ⚠  {tissue}: not found — skipping")
            continue
        img_src = report_dir / "img"
        de_src  = report_dir / "de"
        cards = []

        for safe_dir in sorted(img_src.iterdir()):
            if not safe_dir.is_dir():
                continue
            safe = safe_dir.name
            dst_img = OUTPUT_DIR / "img" / tissue / safe
            dst_img.mkdir(parents=True, exist_ok=True)
            for f in sorted(safe_dir.glob("*.png")):
                shutil.copy2(f, dst_img / f.name)

            def _img(name):
                return f"img/{tissue}/{safe}/{name}" if (dst_img / name).exists() else None

            de_csv_path = None
            de_json = None
            if de_src.exists():
                src_csv = de_src / f"{safe}.csv"
                if src_csv.exists():
                    dst_csv = OUTPUT_DIR / "de" / f"{tissue}_{safe}.csv"
                    shutil.copy2(src_csv, dst_csv)
                    de_csv_path = f"de/{tissue}_{safe}.csv"
                    de_json = pd.read_csv(dst_csv).to_json(orient="split")

            cards.append({
                "label": safe.replace("_", " "),
                "anchor": safe,
                "disease": safe,
                "stats": {},
                "disease_breakdown": [],
                "umap_celltype_path":  _img("umap_celltype.png"),
                "umap_disease_path":   _img("umap_disease.png"),
                "umap_highlight_path": _img("umap_highlight.png"),
                "volcano_path":        _img("volcano.png"),
                "dotplot_path":        None,
                "de_csv_path":         de_csv_path,
                "de_json":             de_json,
            })
        tissues[tissue] = cards
        print(f"  {tissue}: {len(cards)} subclusters")

    total = sum(len(v) for v in tissues.values())
    template_dir = Path(__file__).parent / "templates"
    env = Environment(loader=FileSystemLoader(str(template_dir)), autoescape=False)
    template = env.get_template("combined_report.html.j2")
    html = template.render(tissues=tissues, total_subclusters=total,
                           generated_at=datetime.now().strftime("%Y-%m-%d %H:%M"))
    html_path = OUTPUT_DIR / "index.html"
    html_path.write_text(html, encoding="utf-8")
    print(f"  ✓ Combined HTML: {html_path} ({html_path.stat().st_size // 1024} KB)")


# ─────────────────────────────────────────────────────────────────────────────
# 5. Main
# ─────────────────────────────────────────────────────────────────────────────

def main():
    print(f"\n{'#' * 60}")
    print("# Disease Subcluster Pipeline — End-to-End Smoke Test")
    print(f"{'#' * 60}")

    with tempfile.TemporaryDirectory(prefix="sc_pipeline_test_") as tmp:
        out = Path(tmp)
        print(f"\nUsing temp dir: {out}\n")

        print("Creating synthetic AnnData ...", end=" ")
        adata = make_synthetic_adata()
        print(f"OK  ({adata.n_obs} cells × {adata.n_vars} genes)")

        # NB01
        info_df = run_nb01(adata, out)

        # NB02
        run_nb02(adata, out)

        # NB03 combined
        run_nb03_combined(out)

        # Final check
        expected = [
            out / "test_dataset_disease_subclusters.h5ad",
            out / "test_dataset_subcluster_info.csv",
            out / "test_tissue_subcluster_report" / "index.html",
            out / "combined_report" / "index.html",
        ]
        print("\n" + "=" * 60)
        print("Output check:")
        all_ok = True
        for p in expected:
            ok = p.exists()
            all_ok = all_ok and ok
            mark = "✓" if ok else "✗ MISSING"
            size = f"  ({p.stat().st_size // 1024} KB)" if ok else ""
            print(f"  {mark}  {p.relative_to(out)}{size}")

    if all_ok:
        print("\n✅  All pipeline steps completed successfully.\n")
        sys.exit(0)
    else:
        print("\n❌  Some outputs are missing — see above.\n")
        sys.exit(1)


if __name__ == "__main__":
    main()
